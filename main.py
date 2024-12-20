import sys
from PyPDF2 import PdfReader, PdfWriter
from pptx import Presentation
from PyQt6.QtWidgets import (
    QApplication,
    QMainWindow,
    QWidget,
    QVBoxLayout,
    QHBoxLayout,
    QPushButton,
    QLineEdit,
    QLabel,
    QFileDialog,
    QMessageBox,
    QFrame,
    QStackedWidget,
)
from PyQt6.QtCore import Qt, QThread, pyqtSignal
from PyQt6.QtGui import QIcon
import os
import requests
import json
import zipfile
import shutil
from packaging import version

CURRENT_VERSION = "1.0.1"
GITHUB_REPO = "siam500561/python_splitter"


class UpdateChecker(QThread):
    update_available = pyqtSignal(str, str)  # version, download_url

    def run(self):
        try:
            # Check latest release from GitHub
            response = requests.get(
                f"https://api.github.com/repos/{GITHUB_REPO}/releases/latest"
            )
            if response.status_code == 200:
                release_info = response.json()
                latest_version = release_info["tag_name"].replace("v", "")
                print(latest_version)

                # Compare versions
                if version.parse(latest_version) > version.parse(CURRENT_VERSION):
                    download_url = release_info["assets"][0]["browser_download_url"]
                    self.update_available.emit(latest_version, download_url)
        except Exception as e:
            print(f"Error checking for updates: {str(e)}")


class PDFSplitter(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Splitter Updated")
        self.setFixedSize(500, 450)
        # Set window icon
        self.setWindowIcon(QIcon("logo.ico"))
        self.setStyleSheet(
            """
            QMainWindow {
                background-color: white;
            }
            QLabel {
                color: rgb(102, 102, 102);
                font-size: 13px;
            }
        """
        )

        # Check for updates
        self.check_for_updates()

        # Create central widget and main layout
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout(central_widget)
        main_layout.setContentsMargins(20, 20, 20, 20)
        main_layout.setSpacing(15)

        # Create buttons layout
        buttons_layout = QHBoxLayout()
        buttons_layout.setSpacing(10)

        # Create action buttons with icons
        self.extract_pdf_btn = self.create_action_button("Extract PDF", "📄")
        self.merge_pdfs_btn = self.create_action_button("Merge PDFs", "📑")
        self.extract_slides_btn = self.create_action_button("Extract Slides", "📊")

        buttons_layout.addWidget(self.extract_pdf_btn)
        buttons_layout.addWidget(self.merge_pdfs_btn)
        buttons_layout.addWidget(self.extract_slides_btn)

        main_layout.addLayout(buttons_layout)

        # Create stacked widget for different form layouts
        self.stacked_widget = QStackedWidget()

        # Create different form sections
        self.extract_form = self.create_extract_form()
        self.merge_form = self.create_merge_form()
        self.slides_form = self.create_slides_form()

        self.stacked_widget.addWidget(self.extract_form)
        self.stacked_widget.addWidget(self.merge_form)
        self.stacked_widget.addWidget(self.slides_form)

        main_layout.addWidget(self.stacked_widget)

        # Create process button
        self.process_btn = QPushButton("Process PDF")
        self.process_btn.setFixedHeight(40)
        self.process_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self.process_btn.setStyleSheet(
            """
            QPushButton {
                background-color: #2196F3;
                color: white;
                border-radius: 4px;
                font-size: 14px;
                border: none;
            }
            QPushButton:hover {
                background-color: #1976D2;
            }
        """
        )
        main_layout.addWidget(self.process_btn)

        # Connect signals
        self.extract_pdf_btn.clicked.connect(lambda: self.switch_mode("extract"))
        self.merge_pdfs_btn.clicked.connect(lambda: self.switch_mode("merge"))
        self.extract_slides_btn.clicked.connect(lambda: self.switch_mode("slides"))
        self.process_btn.clicked.connect(self.process_pdf)

        # Set initial mode
        self.current_mode = "extract"
        self.switch_mode("extract")

    def check_for_updates(self):
        self.update_checker = UpdateChecker()
        self.update_checker.update_available.connect(self.show_update_dialog)
        self.update_checker.start()

    def show_update_dialog(self, new_version, download_url):
        reply = QMessageBox.question(
            self,
            "Update Available",
            f"A new version (v{new_version}) is available. Would you like to update?",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
        )

        if reply == QMessageBox.StandardButton.Yes:
            self.download_and_install_update(download_url)

    def download_and_install_update(self, download_url):
        try:
            # Create an 'updates' directory if it doesn't exist
            updates_dir = "updates"
            if not os.path.exists(updates_dir):
                os.makedirs(updates_dir)

            # Download the update
            response = requests.get(download_url, stream=True)
            update_zip = os.path.join(updates_dir, "update.zip")

            with open(update_zip, "wb") as f:
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)

            # Extract update to temp directory
            extract_dir = os.path.join(updates_dir, "temp")
            if os.path.exists(extract_dir):
                shutil.rmtree(extract_dir)

            with zipfile.ZipFile(update_zip, "r") as zip_ref:
                zip_ref.extractall(extract_dir)

            # Find the Splitter folder in the extracted contents
            splitter_folder = None
            for item in os.listdir(extract_dir):
                if item == "Splitter":
                    splitter_folder = os.path.join(extract_dir, item)
                    break

            if not splitter_folder:
                raise Exception("Invalid update package: Splitter folder not found")

            # Create a batch file to perform the update
            batch_script = os.path.join(updates_dir, "update.bat")
            current_dir = os.path.abspath(os.path.dirname(__file__))

            with open(batch_script, "w") as f:
                f.write("@echo off\n")
                f.write("timeout /t 1 /nobreak >nul\n")  # Wait for the app to close

                # Copy new files
                f.write(f'xcopy /s /y "{splitter_folder}\\*" "{current_dir}\\"\n')

                # Clean up
                f.write(f'rmdir /s /q "{updates_dir}"\n')

                # Start the updated application
                f.write(f'start "" "{current_dir}\\Splitter.exe"\n')

                # Delete the batch file itself
                f.write('(goto) 2>nul & del "%~f0"')

            # Show success message and start the update process
            QMessageBox.information(
                self,
                "Update Ready",
                "The update has been downloaded. The application will now close and update itself.",
            )

            # Start the batch file and close the application
            os.startfile(batch_script)
            sys.exit(0)

        except Exception as e:
            QMessageBox.critical(
                self,
                "Update Failed",
                f"Failed to install update: {str(e)}\nPlease try again or download the new version manually.",
            )
            # Clean up
            if os.path.exists(updates_dir):
                try:
                    shutil.rmtree(updates_dir)
                except:
                    pass

    def create_action_button(self, text, icon_text):
        btn = QPushButton()
        btn.setFixedSize(150, 90)
        btn.setCursor(Qt.CursorShape.PointingHandCursor)

        # Create vertical layout for icon and text
        layout = QVBoxLayout(btn)
        layout.setAlignment(Qt.AlignmentFlag.AlignCenter)
        layout.setSpacing(5)
        layout.setContentsMargins(10, 10, 10, 10)

        # Add icon label
        icon_label = QLabel(icon_text)
        icon_label.setStyleSheet("font-size: 28px; color: #444;")
        icon_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        layout.addWidget(icon_label)

        # Add text label
        text_label = QLabel(text)
        text_label.setStyleSheet("color: #666; font-size: 13px;")
        text_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        layout.addWidget(text_label)

        btn.setStyleSheet(
            """
            QPushButton {
                background-color: white;
                border: 1px solid #e8e8e8;
                border-radius: 8px;
            }
            QPushButton:hover {
                background-color: #f5f5f5;
                border: 1px solid #e0e0e0;
            }
        """
        )
        return btn

    def create_styled_input(self, placeholder, is_readonly=True):
        input_field = QLineEdit()
        input_field.setPlaceholderText(placeholder)
        input_field.setReadOnly(is_readonly)
        input_field.setFixedHeight(38)
        if is_readonly:
            input_field.setCursor(Qt.CursorShape.PointingHandCursor)
        input_field.setStyleSheet(
            """
            QLineEdit {
                border: 1px solid #e0e0e0;
                border-radius: 4px;
                padding: 8px;
                background-color: white;
                color: #444;
                font-size: 13px;
            }
            QLineEdit:disabled {
                background-color: #f5f5f5;
            }
        """
        )
        if is_readonly:
            input_field.mousePressEvent = lambda _: self.select_file(input_field)
        return input_field

    def create_extract_form(self):
        form = QFrame()
        layout = QVBoxLayout(form)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(4)

        # Input File
        input_label = QLabel("Input File")
        input_label.setContentsMargins(0, 0, 0, 0)
        layout.addWidget(input_label)
        self.extract_input_file = self.create_styled_input("Select PDF file...")
        layout.addWidget(self.extract_input_file)

        # Page Range
        range_label = QLabel("Page Range")
        range_label.setContentsMargins(0, 6, 0, 0)
        layout.addWidget(range_label)
        self.page_range = self.create_styled_input("e.g., 1,2 or 1-3,5-7", False)
        # Connect textChanged signal to update output filename
        self.page_range.textChanged.connect(self.update_extract_output_filename)
        layout.addWidget(self.page_range)

        # Output File
        output_label = QLabel("Output File")
        output_label.setContentsMargins(0, 6, 0, 0)
        layout.addWidget(output_label)
        self.extract_output_file = self.create_styled_input("Output file path...")
        layout.addWidget(self.extract_output_file)

        return form

    def create_merge_form(self):
        form = QFrame()
        layout = QVBoxLayout(form)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(4)

        # First Input File
        first_label = QLabel("First PDF File")
        first_label.setContentsMargins(0, 0, 0, 0)
        layout.addWidget(first_label)
        self.merge_input_file1 = self.create_styled_input("Select first PDF file...")
        layout.addWidget(self.merge_input_file1)

        # Second Input File
        second_label = QLabel("Second PDF File")
        second_label.setContentsMargins(0, 6, 0, 0)
        layout.addWidget(second_label)
        self.merge_input_file2 = self.create_styled_input("Select second PDF file...")
        layout.addWidget(self.merge_input_file2)

        # Output File
        output_label = QLabel("Output File")
        output_label.setContentsMargins(0, 6, 0, 0)
        layout.addWidget(output_label)
        self.merge_output_file = self.create_styled_input("Output file path...")
        layout.addWidget(self.merge_output_file)

        return form

    def create_slides_form(self):
        form = QFrame()
        layout = QVBoxLayout(form)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(4)

        # Input File
        input_label = QLabel("Input File")
        input_label.setContentsMargins(0, 0, 0, 0)
        layout.addWidget(input_label)
        self.slides_input_file = self.create_styled_input("Select PPTX file...")
        layout.addWidget(self.slides_input_file)

        # Slide Range
        range_label = QLabel("Slide Range")
        range_label.setContentsMargins(0, 6, 0, 0)
        layout.addWidget(range_label)
        self.slides_range = self.create_styled_input("e.g., 1,2 or 1-3,5-7", False)
        # Connect textChanged signal to update output filename
        self.slides_range.textChanged.connect(self.update_slides_output_filename)
        layout.addWidget(self.slides_range)

        # Output File
        output_label = QLabel("Output File")
        output_label.setContentsMargins(0, 6, 0, 0)
        layout.addWidget(output_label)
        self.slides_output_file = self.create_styled_input("Output file path...")
        layout.addWidget(self.slides_output_file)

        return form

    def select_file(self, input_field):
        if "PPTX" in input_field.placeholderText():
            file_name, _ = QFileDialog.getOpenFileName(
                self, "Select PPTX file", "", "PowerPoint files (*.pptx)"
            )
        else:
            file_name, _ = QFileDialog.getOpenFileName(
                self, "Select PDF file", "", "PDF files (*.pdf)"
            )

        if file_name:
            input_field.setText(file_name)
            # Set default output name
            dir_path = os.path.dirname(file_name)
            base_name = os.path.splitext(os.path.basename(file_name))[0]

            if self.current_mode == "extract":
                # Don't set output file name here, it will be set when page range is entered
                self.extract_output_file.clear()
            elif self.current_mode == "merge":
                self.merge_output_file.setText(
                    os.path.join(dir_path, f"{base_name}_merged.pdf")
                )
            elif self.current_mode == "slides":
                self.slides_output_file.setText(
                    os.path.join(dir_path, f"{base_name}_slides.pdf")
                )

    def update_extract_output_filename(self):
        if self.extract_input_file.text() and self.page_range.text():
            dir_path = os.path.dirname(self.extract_input_file.text())
            base_name = os.path.splitext(
                os.path.basename(self.extract_input_file.text())
            )[0]
            # Clean up page range for filename
            page_range = self.page_range.text().replace(" ", "").replace(",", "-")
            self.extract_output_file.setText(
                os.path.join(dir_path, f"{base_name}_pages_{page_range}.pdf")
            )

    def switch_mode(self, mode):
        self.current_mode = mode

        # Reset button colors to white with lighter border
        default_style = """
            QPushButton {
                background-color: white;
                border: 1px solid #f0f0f0;
                border-radius: 8px;
                cursor: pointer;
            }
            QPushButton:hover {
                background-color: #f5f5f5;
                border: 1px solid #e0e0e0;
            }
        """

        self.extract_pdf_btn.setStyleSheet(default_style)
        self.merge_pdfs_btn.setStyleSheet(default_style)
        self.extract_slides_btn.setStyleSheet(default_style)

        # Highlight selected button with a subtle blue
        active_style = """
            QPushButton {
                background-color: #f8f9fa;
                border: 1px solid #2196F3;
                border-radius: 8px;
                cursor: pointer;
            }
            QPushButton:hover {
                background-color: #f8f9fa;
                border: 1px solid #2196F3;
            }
        """

        # Switch form and button style based on mode
        if mode == "extract":
            self.extract_pdf_btn.setStyleSheet(active_style)
            self.stacked_widget.setCurrentIndex(0)
            self.process_btn.setText("Extract PDF")
        elif mode == "merge":
            self.merge_pdfs_btn.setStyleSheet(active_style)
            self.stacked_widget.setCurrentIndex(1)
            self.process_btn.setText("Merge PDFs")
        else:  # slides
            self.extract_slides_btn.setStyleSheet(active_style)
            self.stacked_widget.setCurrentIndex(2)
            self.process_btn.setText("Extract Slides")

    def process_pdf(self):
        if self.current_mode == "extract":
            self.extract_pages()
        elif self.current_mode == "merge":
            self.merge_pdfs()
        else:
            self.extract_slides()

    def open_file_in_explorer(self, file_path):
        try:
            # Convert to absolute path and normalize
            abs_path = os.path.abspath(file_path)

            if sys.platform == "win32":
                # Use double quotes and replace forward slashes with backslashes
                normalized_path = abs_path.replace("/", "\\")
                os.system(f'explorer /select,"{normalized_path}"')
            elif sys.platform == "darwin":  # macOS
                os.system(f'open -R "{abs_path}"')
            else:  # Linux
                os.system(f'xdg-open "{os.path.dirname(abs_path)}"')
        except Exception as e:
            print(f"Error opening file location: {str(e)}")

    def extract_pages(self):
        input_file = self.extract_input_file.text()
        output_file = self.extract_output_file.text()
        page_range = self.page_range.text()

        if not input_file or not output_file or not page_range:
            QMessageBox.warning(self, "Error", "Please fill in all fields.")
            return

        try:
            reader = PdfReader(input_file)
            total_pages = len(reader)
            writer = PdfWriter()

            # Parse page range
            pages = []
            for part in page_range.split(","):
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    if start > end:
                        QMessageBox.warning(
                            self,
                            "Error",
                            "Invalid page range: start page cannot be greater than end page.",
                        )
                        return
                    if start < 1 or end > total_pages:
                        QMessageBox.warning(
                            self,
                            "Error",
                            f"Page range {start}-{end} is out of bounds. PDF has {total_pages} pages.",
                        )
                        return
                    pages.extend(range(start - 1, end))
                else:
                    page = int(part)
                    if page < 1 or page > total_pages:
                        QMessageBox.warning(
                            self,
                            "Error",
                            f"Page {page} is out of bounds. PDF has {total_pages} pages.",
                        )
                        return
                    pages.append(page - 1)

            # Add pages to writer
            for page_num in pages:
                writer.add_page(reader.pages[page_num])

            # Save the output file
            with open(output_file, "wb") as output:
                writer.write(output)

            # Open the file location in explorer
            os.system(f'explorer /select,"{os.path.abspath(output_file)}"')

        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred: {str(e)}")
            return

    def merge_pdfs(self):
        try:
            if not self.merge_input_file1.text() or not self.merge_input_file2.text():
                QMessageBox.warning(self, "Error", "Please select both input files.")
                return

            if not self.merge_output_file.text():
                QMessageBox.warning(self, "Error", "Please select an output file.")
                return

            writer = PdfWriter()

            # Add pages from first PDF
            reader1 = PdfReader(self.merge_input_file1.text())
            for page in reader1.pages:
                writer.add_page(page)

            # Add pages from second PDF
            reader2 = PdfReader(self.merge_input_file2.text())
            for page in reader2.pages:
                writer.add_page(page)

            # Write the merged PDF
            output_path = self.merge_output_file.text()
            with open(output_path, "wb") as output_file:
                writer.write(output_file)

            # Open the file location
            self.open_file_in_explorer(output_path)

        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred: {str(e)}")

    def extract_slides(self):
        try:
            if not self.slides_input_file.text():
                QMessageBox.warning(self, "Error", "Please select a PPTX file.")
                return

            if not self.slides_output_file.text():
                QMessageBox.warning(self, "Error", "Please select an output file.")
                return

            if not self.slides_range.text():
                QMessageBox.warning(self, "Error", "Please enter slide range.")
                return

            # Load source presentation
            source_prs = Presentation(self.slides_input_file.text())
            total_slides = len(source_prs.slides)

            # Parse slide range
            slides_to_extract = self.parse_page_range(
                self.slides_range.text(), total_slides
            )

            # Generate base output path
            file_path = self.slides_output_file.text()
            dir_path = os.path.dirname(file_path)
            base_name = os.path.splitext(os.path.basename(file_path))[0]

            # Create range string for filename
            range_str = "_".join(
                part.strip() for part in self.slides_range.text().split(",")
            )
            range_str = range_str.replace("-", "_")
            output_path = os.path.join(dir_path, f"{base_name}_slides_{range_str}.pptx")

            # Create new presentation by loading the source
            dest_prs = Presentation(self.slides_input_file.text())

            # Keep only the slides we want
            slides_to_keep = [
                dest_prs.slides._sldIdLst[idx] for idx in slides_to_extract
            ]

            # Remove all other slides
            for slide in dest_prs.slides._sldIdLst[:]:
                if slide not in slides_to_keep:
                    rId = slide.rId
                    dest_prs.part.drop_rel(rId)
                    dest_prs.slides._sldIdLst.remove(slide)

            # Save the presentation with selected slides
            dest_prs.save(output_path)

            # Open the exact output file in explorer
            self.open_file_in_explorer(output_path)

        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred: {str(e)}")

    def parse_page_range(self, range_str, max_pages):
        pages = set()
        parts = range_str.split(",")

        for part in parts:
            if "-" in part:
                start, end = map(int, part.split("-"))
                pages.update(range(start - 1, end))
            else:
                pages.add(int(part) - 1)

        return sorted(list(pages))

    def update_slides_output_filename(self):
        if self.slides_input_file.text() and self.slides_range.text():
            dir_path = os.path.dirname(self.slides_input_file.text())
            base_name = os.path.splitext(
                os.path.basename(self.slides_input_file.text())
            )[0]
            # Clean up range for filename
            range_str = "_".join(
                part.strip() for part in self.slides_range.text().split(",")
            )
            range_str = range_str.replace("-", "_")
            self.slides_output_file.setText(
                os.path.join(dir_path, f"{base_name}_slides_{range_str}.pptx")
            )


def main():
    app = QApplication(sys.argv)
    window = PDFSplitter()
    window.show()
    sys.exit(app.exec())


if __name__ == "__main__":
    main()
